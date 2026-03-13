"""Document loading and chunking utilities for text and multimodal artifacts."""

from __future__ import annotations

import mimetypes
from dataclasses import dataclass
from pathlib import Path
from typing import Type

import fitz
from langchain_core.document_loaders import BaseLoader
from langchain_core.documents import Document
from langchain_community.document_loaders import (
    CSVLoader,
    PyMuPDFLoader,
    TextLoader,
    UnstructuredExcelLoader,
    UnstructuredMarkdownLoader,
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.config import CHUNK_OVERLAP, CHUNK_SIZE, DOCS_DIR, SUPPORTED_EXTENSIONS, VISION_MAX_IMAGES_PER_DOC


def safe_print(msg: str) -> None:
    """Print with fallback for encoding issues on Windows."""
    try:
        print(msg)
    except UnicodeEncodeError:
        # Replace problematic characters with ?
        safe_msg = msg.encode("ascii", errors="replace").decode("ascii")
        print(safe_msg)


@dataclass
class VisualArtifact:
    """Normalized image artifact extracted from documents for captioning."""

    content: bytes
    metadata: dict


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".tiff"}


def load_pdf_text(file_path: Path) -> list[Document]:
    """Load text from a PDF using LangChain's PyMuPDF loader."""
    loader = PyMuPDFLoader(str(file_path))
    return loader.load()


def load_pptx_text(file_path: Path) -> list[Document]:
    """Load PowerPoint slide text using python-pptx directly."""
    prs = Presentation(str(file_path))
    text_parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))

    content = "\n\n".join(text_parts)
    if not content.strip():
        return []
    return [Document(page_content=content, metadata={"source": str(file_path)})]


def load_docx(file_path: str) -> list[Document]:
    """Load Word documents using python-docx directly."""
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                paragraphs.append(" | ".join(row_text))

    content = "\n\n".join(paragraphs)
    if not content.strip():
        return []

    return [Document(page_content=content, metadata={"source": file_path})]


def extract_pdf_images(file_path: Path, max_images: int = VISION_MAX_IMAGES_PER_DOC) -> list[VisualArtifact]:
    """Extract image bytes from a PDF with page-aware metadata."""
    artifacts: list[VisualArtifact] = []
    doc = fitz.open(str(file_path))
    try:
        for page_index in range(len(doc)):
            page = doc[page_index]
            images = page.get_images(full=True)
            for image_index, image_info in enumerate(images, 1):
                if len(artifacts) >= max_images:
                    return artifacts
                xref = image_info[0]
                base = doc.extract_image(xref)
                image_bytes = base.get("image")
                if not image_bytes:
                    continue
                ext = (base.get("ext") or "").strip().lower()
                mime_type = mimetypes.types_map.get(f".{ext}", "application/octet-stream")
                artifacts.append(
                    VisualArtifact(
                        content=image_bytes,
                        metadata={
                            "source": str(file_path),
                            "file_type": "pdf",
                            "page_or_slide": page_index + 1,
                            "image_index": image_index,
                            "image_ext": ext,
                            "image_mime": mime_type,
                            "modality": "image",
                        },
                    )
                )
    finally:
        doc.close()
    return artifacts


def extract_pptx_images(file_path: Path, max_images: int = VISION_MAX_IMAGES_PER_DOC) -> list[VisualArtifact]:
    """Extract image bytes from PPTX picture shapes with slide-aware metadata."""
    artifacts: list[VisualArtifact] = []
    prs = Presentation(str(file_path))
    for slide_index, slide in enumerate(prs.slides, 1):
        slide_image_index = 0
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            if len(artifacts) >= max_images:
                return artifacts
            image = shape.image
            if image is None:
                continue
            slide_image_index += 1
            ext = (image.ext or "").strip().lower()
            mime_type = image.content_type or mimetypes.types_map.get(f".{ext}", "application/octet-stream")
            artifacts.append(
                VisualArtifact(
                    content=image.blob,
                    metadata={
                        "source": str(file_path),
                        "file_type": "pptx",
                        "page_or_slide": slide_index,
                        "image_index": slide_image_index,
                        "image_ext": ext,
                        "image_mime": mime_type,
                        "modality": "image",
                    },
                )
            )
    return artifacts


def load_image_file(file_path: Path) -> list[VisualArtifact]:
    """Load a standalone image as a visual artifact for vision captioning."""
    ext = file_path.suffix.lower().lstrip(".")
    mime_type = mimetypes.types_map.get(file_path.suffix.lower(), "application/octet-stream")
    return [
        VisualArtifact(
            content=file_path.read_bytes(),
            metadata={
                "source": str(file_path),
                "file_type": file_path.suffix.lower().lstrip("."),
                "page_or_slide": 1,
                "image_index": 1,
                "image_ext": ext,
                "image_mime": mime_type,
                "modality": "image",
            },
        )
    ]


# Map file extensions to their respective loaders (text-oriented flow).
LOADER_MAP: dict[str, Type[BaseLoader] | str] = {
    ".txt": TextLoader,
    ".md": UnstructuredMarkdownLoader,
    ".pdf": "pdf_text",
    ".docx": "docx",  # Custom loader
    ".xlsx": UnstructuredExcelLoader,
    ".xls": UnstructuredExcelLoader,
    ".csv": CSVLoader,
    ".pptx": "pptx_text",  # Custom loader
}

# Legacy formats that need LibreOffice - we'll skip these with a warning
LEGACY_FORMATS = {".doc", ".ppt"}


def get_loader_for_file(file_path: Path) -> BaseLoader | None:
    """Get the appropriate loader for a file based on its extension."""
    ext = file_path.suffix.lower()
    
    if ext in LEGACY_FORMATS:
        return None  # Skip legacy formats
    
    if ext not in LOADER_MAP:
        return None

    loader_class = LOADER_MAP[ext]

    # Custom loaders return string markers
    if isinstance(loader_class, str):
        return loader_class  # type: ignore

    # CSVLoader and TextLoader need encoding specified
    if loader_class in (CSVLoader, TextLoader):
        return loader_class(str(file_path), encoding="utf-8")

    return loader_class(str(file_path))


def load_documents(docs_dir: Path = DOCS_DIR) -> list[Document]:
    """Load all supported documents from the specified directory (recursive)."""
    documents: list[Document] = []
    skipped_legacy: list[str] = []

    if not docs_dir.exists():
        raise FileNotFoundError(f"Documents directory not found: {docs_dir}")

    # Collect all files with supported extensions (recursive search)
    files_to_load: list[Path] = []
    for ext in SUPPORTED_EXTENSIONS:
        files_to_load.extend(docs_dir.glob(f"**/*{ext}"))

    # Filter out temp files (e.g., ~$document.docx)
    files_to_load = [f for f in files_to_load if not f.name.startswith("~$")]

    # Sort for consistent ordering
    files_to_load = sorted(files_to_load, key=lambda x: str(x).lower())

    if not files_to_load:
        return documents

    safe_print(f"Found {len(files_to_load)} files to process...")

    for i, file_path in enumerate(files_to_load, 1):
        ext = file_path.suffix.lower()

        # Track skipped legacy files
        if ext in LEGACY_FORMATS:
            skipped_legacy.append(file_path.name)
            continue

        if ext in IMAGE_EXTENSIONS:
            # Standalone images are handled by multimodal artifact flow.
            continue

        try:
            safe_print(f"  [{i}/{len(files_to_load)}] Loading: {file_path.name}")

            # Custom loaders
            if ext == ".docx":
                loaded_docs = load_docx(str(file_path))
            elif ext == ".pdf":
                loaded_docs = load_pdf_text(file_path)
            elif ext == ".pptx":
                loaded_docs = load_pptx_text(file_path)
            else:
                loader = get_loader_for_file(file_path)
                if loader is None:
                    continue
                loaded_docs = loader.load()

            # Add rich metadata for better context and filtering
            rel_path = file_path.relative_to(docs_dir)
            title = file_path.stem.lower()  # filename without extension, lowercase for filtering
            file_type = file_path.suffix.lower().lstrip(".")  # extension without dot

            for doc in loaded_docs:
                doc.metadata["source"] = str(rel_path)
                doc.metadata["title"] = title
                doc.metadata["file_type"] = file_type
                doc.metadata["modality"] = "text"
            documents.extend(loaded_docs)

        except Exception as e:
            # Log error but continue with other files
            safe_print(f"  Warning: Failed to load {file_path.name}: {e}")

    if skipped_legacy:
        safe_print(f"\nSkipped {len(skipped_legacy)} legacy files (.doc/.ppt) - need LibreOffice to convert")

    return documents


def load_multimodal_artifacts(docs_dir: Path = DOCS_DIR) -> tuple[list[Document], list[VisualArtifact]]:
    """Load text documents and image artifacts from supported local files."""
    text_documents: list[Document] = []
    image_artifacts: list[VisualArtifact] = []

    if not docs_dir.exists():
        raise FileNotFoundError(f"Documents directory not found: {docs_dir}")

    files_to_load: list[Path] = []
    for ext in SUPPORTED_EXTENSIONS:
        files_to_load.extend(docs_dir.glob(f"**/*{ext}"))
    files_to_load = [f for f in files_to_load if not f.name.startswith("~$")]
    files_to_load = sorted(files_to_load, key=lambda x: str(x).lower())

    for file_path in files_to_load:
        ext = file_path.suffix.lower()
        rel_path = str(file_path.relative_to(docs_dir))
        title = file_path.stem.lower()
        file_type = ext.lstrip(".")

        try:
            loaded_text_docs: list[Document] = []
            if ext == ".pdf":
                loaded_text_docs = load_pdf_text(file_path)
                image_artifacts.extend(extract_pdf_images(file_path))
            elif ext == ".pptx":
                loaded_text_docs = load_pptx_text(file_path)
                image_artifacts.extend(extract_pptx_images(file_path))
            elif ext == ".docx":
                loaded_text_docs = load_docx(str(file_path))
            elif ext in IMAGE_EXTENSIONS:
                image_artifacts.extend(load_image_file(file_path))
            else:
                loader = get_loader_for_file(file_path)
                if loader is not None and not isinstance(loader, str):
                    loaded_text_docs = loader.load()

            for doc in loaded_text_docs:
                doc.metadata["source"] = rel_path
                doc.metadata["title"] = title
                doc.metadata["file_type"] = file_type
                doc.metadata["modality"] = "text"
                text_documents.append(doc)

            for artifact in image_artifacts:
                source = artifact.metadata.get("source")
                if source == str(file_path):
                    artifact.metadata["source"] = rel_path
                    artifact.metadata["parent_source"] = rel_path
                    artifact.metadata["title"] = title

        except Exception as e:
            safe_print(f"  Warning: Failed to process {file_path.name}: {e}")

    return text_documents, image_artifacts


def chunk_documents(
    documents: list[Document],
    chunk_size: int = CHUNK_SIZE,
    chunk_overlap: int = CHUNK_OVERLAP
) -> list[Document]:
    """Split documents into smaller chunks for embedding."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""]
    )
    
    chunks = text_splitter.split_documents(documents)
    return chunks


def load_and_chunk_documents(docs_dir: Path = DOCS_DIR) -> list[Document]:
    """Load documents and split them into chunks."""
    documents = load_documents(docs_dir)
    chunks = chunk_documents(documents)
    return chunks
